"""Admin helper: dump every piece of text in a PowerPoint template.

Usage:
    py inspect_template.py "HO3 Proposal 01july2026.pptx"

Writes <template name>_text.txt next to the template, listing every
textbox paragraph and table cell, slide by slide. Use the output to fill
in the exact placeholder keys in a mapping JSON (or paste it back to
your AI assistant and ask it to finish the mapping for you).
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def main() -> None:
    if len(sys.argv) < 2:
        print('Usage: py inspect_template.py "Template file.pptx"')
        sys.exit(1)

    path = Path(sys.argv[1])
    prs = Presentation(path)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n========== SLIDE {i} ==========")
        for shape in iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(f"[TEXTBOX  name={shape.name!r}]")
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        # repr() so invisible characters/line breaks are visible
                        lines.append(f"    {p.text!r}")
            elif shape.has_table:
                lines.append(f"[TABLE    name={shape.name!r}]")
                for r_i, row in enumerate(shape.table.rows):
                    cells = [c.text for c in row.cells]
                    lines.append(f"    row {r_i}: {cells!r}")

    out_path = path.with_name(path.stem + "_text.txt")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    print("\n".join(lines))
    print(f"\nSaved to: {out_path}")


if __name__ == "__main__":
    main()
