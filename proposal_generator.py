"""Quote/proposal generation engine.

All PowerPoint manipulation logic is preserved unchanged from the original
generate_proposal.py. The only structural change: instead of looping over
every row in the spreadsheet, `generate_quote()` accepts a single row and
produces a single .pptx file.
"""

from __future__ import annotations

import re
import sys
from copy import deepcopy
from pathlib import Path
from typing import Any, List, Tuple

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
TEMPLATE_FILENAME = "Proposal  PERSONAL AUTO 23mar2026.pptx"
OUTPUT_DIR_NAME = "Quotes"
SLIDE_REPLACE_LIMIT = 5


class QuoteGenerationError(Exception):
    """Raised when a quote cannot be generated (missing template, etc.)."""


def app_base_dir() -> Path:
    """Folder containing the app.

    Works both when running as a plain script and when frozen into a
    PyInstaller .exe (where __file__ points inside the bundle).
    """
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


def default_template_path(profile_key: str = "AUTO") -> Path:
    """Find the profile's template automatically — the user is never asked.

    Search order (per quote type, matched by filename pattern so renamed
    or re-dated templates still work):
    1. Next to the .exe / script (lets the agency swap templates without
       rebuilding the app)
    2. Inside the PyInstaller bundle (templates embedded via --add-data)
    """
    import quote_profiles
    pattern = quote_profiles.PROFILES[profile_key]["template_glob"]

    search_dirs = [app_base_dir()]
    # sys._MEIPASS is the temp folder PyInstaller unpacks bundled data into
    bundle_dir = getattr(sys, "_MEIPASS", None)
    if bundle_dir:
        search_dirs.append(Path(bundle_dir))

    for folder in search_dirs:
        matches = sorted(folder.glob(pattern))
        if matches:
            return matches[0]
    # Not found; return a representative path so the error message is clear
    return app_base_dir() / pattern.replace("*", "")


def _desktop_dir() -> Path:
    """Locate the user's Desktop folder (Windows-API first, then fallbacks)."""
    if sys.platform == "win32":
        try:
            import ctypes
            buf = ctypes.create_unicode_buffer(260)
            # CSIDL 0x10 = Desktop directory (also works with OneDrive desktops)
            ctypes.windll.shell32.SHGetFolderPathW(None, 0x10, None, 0, buf)
            if buf.value:
                return Path(buf.value)
        except Exception:
            pass
    candidate = Path.home() / "Desktop"
    return candidate if candidate.exists() else Path.home()


def default_output_dir() -> Path:
    """'Quotes' folder on the user's Desktop; created automatically."""
    out = _desktop_dir() / OUTPUT_DIR_NAME
    out.mkdir(parents=True, exist_ok=True)
    return out


# ---------------------------------------------------------------------------
# Value normalisation helpers (unchanged from original script)
# ---------------------------------------------------------------------------
def normalize_value(value: Any) -> str:
    if pd.isna(value):
        return ""
    if isinstance(value, pd.Timestamp):
        return str(value.date())
    return str(value).strip()


def normalize_value_or_na(value: Any) -> str:
    normalized = normalize_value(value)
    return normalized if normalized else "N/A"


def format_date_for_filename(value: Any) -> str:
    if pd.isna(value):
        return "UnknownDate"
    try:
        ts = pd.to_datetime(value)
        return ts.strftime("%Y%m%d")
    except (ValueError, TypeError):
        cleaned = re.sub(r"[^A-Za-z0-9_-]", "", str(value))
        return cleaned or "UnknownDate"


def sanitize_filename_part(value: str) -> str:
    cleaned = re.sub(r'[<>:"/|?*]', "", value)
    cleaned = re.sub(r"\s+", "_", cleaned).strip("_ ")
    return cleaned or "Client"


def get_named_insured_lines(named_insureds: str) -> List[str]:
    return [part.strip() for part in named_insureds.split(",") if part.strip()]


# ---------------------------------------------------------------------------
# Placeholder mapping (unchanged from original script)
# ---------------------------------------------------------------------------
def build_text_mapping(row: pd.Series) -> Tuple[dict, List[str], str]:
    client_name = normalize_value(row["Client Name"])
    proposal_date = normalize_value(row["Proposal Date"])
    lead_producer = normalize_value(row["Lead Producer"])
    secondary_producer = normalize_value(row["Secondary Producer"])
    if lead_producer and secondary_producer:
        producer_text = f"{lead_producer} / {secondary_producer}"
    else:
        producer_text = lead_producer or secondary_producer

    named_insureds_raw = normalize_value(row["Named Insureds"])
    named_insured_lines = get_named_insured_lines(named_insureds_raw)
    named_insureds_text = "\n".join(named_insured_lines)

    address = ", ".join(
        part for part in [
            normalize_value(row["Street"]),
            normalize_value(row["City"]),
            normalize_value(row["State"]),
            normalize_value(row["ZIP"]),
        ]
        if part
    )

    mapping = {
        "SUSAN LOMAN": client_name,
        "TEST CO #2": client_name,
        "March 21, 2026": proposal_date.upper(),
        "Lillian Wilson, Broker": producer_text.upper(),
        "Susan Loman\nRichard Loman": named_insureds_text,
        "2821 N Ocean Blvd Fort Lauderdale Florida 33308": address,
        "Expiring carrier": normalize_value(row["Expiring Carrier"]),
        "Limit - BI": normalize_value(row["Current BI Limit"]),
        "Limit - PD": normalize_value(row["Current PD Limit"]),
        "Limit - UIM": normalize_value(row["Current UIM Limit"]),
        "Limit - MED": normalize_value(row["Current MED Limit"]),
        "Limit -TE": normalize_value(row["Current TE Limit"]),
        "Limit -PIP": normalize_value(row["Current PIP Limit"]),
        "DED - Col": normalize_value(row["Current Collision Deductible"]),
        "DEC-COMP": normalize_value(row["Current Comp Deductible"]),
        "Old Prem": normalize_value(row["Current Annual Premium"]),
        "New carrier": normalize_value(row["New Carrier"]),
        "new- BI": normalize_value(row["New Renewal BI Limit"]),
        "new- PD": normalize_value(row["New Renewal PD Limit"]),
        "new - UIM": normalize_value(row["New Renewal UIM Limit"]),
        "new - MED": normalize_value(row["New Renewal MED Limit"]),
        "new-TE": normalize_value(row["New Renewal TE Limit"]),
        "new -PIP": normalize_value(row["New PIP Limit"]),
        "new - Col": normalize_value(row["New Collision Deductible"]),
        "new-COMP": normalize_value(row["New Comp Deductible"]),
        "new -prem": normalize_value(row["New Premium"]),
        "LOB": normalize_value(row["LOB"]).upper(),
        "Ambest": normalize_value(row["A.M. Best Rating"]),
        "Commercial General Liability": normalize_value(row["LOB"]).upper(),
    }

    return mapping, named_insured_lines, client_name


# ---------------------------------------------------------------------------
# PowerPoint shape/text helpers (unchanged from original script)
# ---------------------------------------------------------------------------
def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def replace_in_runs(paragraph, mapping: dict) -> None:
    for run in paragraph.runs:
        text = run.text or ""
        if not text:
            continue
        new_text = text
        for placeholder in sorted(mapping.keys(), key=len, reverse=True):
            if placeholder and placeholder in new_text:
                new_text = new_text.replace(placeholder, mapping[placeholder])
        if new_text != text:
            run.text = new_text


def set_cell_text(cell, new_text: str) -> None:
    tf = cell.text_frame
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        cell.text = new_text
        return
    first_p = paragraphs[0]
    runs = list(first_p.runs)
    if runs:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
    else:
        first_p.text = new_text
    txBody = tf._txBody
    for extra in paragraphs[1:]:
        txBody.remove(extra._p)


def replace_text_on_slide(slide, mapping: dict) -> None:
    for shape in iter_shapes(slide.shapes):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_in_runs(paragraph, mapping)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        replace_in_runs(paragraph, mapping)


def fill_slide2_named_insureds_textbox(slide, insured_lines: List[str]) -> None:
    if not insured_lines:
        return
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        has_susan = any("Susan Loman" in p.text for p in tf.paragraphs)
        has_richard = any("Richard Loman" in p.text for p in tf.paragraphs)
        if not (has_susan and has_richard):
            continue

        template_p = deepcopy(tf.paragraphs[0]._p)
        txBody = tf._txBody
        for paragraph in list(tf.paragraphs):
            txBody.remove(paragraph._p)

        a_r = qn("a:r")
        a_t = qn("a:t")
        for name in insured_lines:
            new_p = deepcopy(template_p)
            runs = new_p.findall(a_r)
            if runs:
                t_elem = runs[0].find(a_t)
                if t_elem is None:
                    t_elem = runs[0].makeelement(a_t, {})
                    runs[0].append(t_elem)
                t_elem.text = name
                for extra in runs[1:]:
                    new_p.remove(extra)
            else:
                r = new_p.makeelement(a_r, {})
                t = new_p.makeelement(a_t, {})
                t.text = name
                r.append(t)
                new_p.append(r)
            txBody.append(new_p)
        return


def fill_driver_table(slide, row_data: pd.Series) -> None:
    for shape in iter_shapes(slide.shapes):
        if not shape.has_table:
            continue
        table = shape.table
        header_text = " ".join(cell.text for cell in table.rows[0].cells)
        if "Name" not in header_text:
            continue
        if "VIN" in header_text:
            continue

        drivers = [
            (
                normalize_value(row_data["Driver 1 Name"]),
                normalize_value(row_data["Driver 1 DL#"]),
                normalize_value(row_data["Driver 1 State"]),
            ),
            (
                normalize_value_or_na(row_data["Driver 2 Name"]),
                normalize_value_or_na(row_data["Driver 2 DL#"]),
                normalize_value_or_na(row_data["Driver 2 State"]),
            ),
        ]

        for i, (name, dl, state) in enumerate(drivers, start=1):
            if i >= len(table.rows):
                break
            row = table.rows[i]
            set_cell_text(row.cells[1], name)
            set_cell_text(row.cells[2], dl)
            set_cell_text(row.cells[3], state)
        return


def fill_vehicle_table(slide, row_data: pd.Series) -> None:
    for shape in iter_shapes(slide.shapes):
        if not shape.has_table:
            continue
        table = shape.table
        header_text = " ".join(cell.text for cell in table.rows[0].cells)
        if "VIN" not in header_text:
            continue

        vehicles = [
            (
                normalize_value(row_data["Vehicle 1 VIN"]),
                normalize_value(row_data["Vehicle 1 Make"]),
                normalize_value(row_data["Vehicle 1 Model"]),
                normalize_value(row_data["Vehicle 1 Year"]),
            ),
            (
                normalize_value_or_na(row_data["Vehicle 2 VIN"]),
                normalize_value_or_na(row_data["Vehicle 2 Make"]),
                normalize_value_or_na(row_data["Vehicle 2 Model"]),
                normalize_value_or_na(row_data["Vehicle 2 Year"]),
            ),
        ]

        for i, (vin, make, model, year) in enumerate(vehicles, start=1):
            if i >= len(table.rows):
                break
            row = table.rows[i]
            set_cell_text(row.cells[1], vin)
            set_cell_text(row.cells[2], make)
            set_cell_text(row.cells[3], model)
            set_cell_text(row.cells[4], year)
        return


def fill_slide3_address_table(slide, address: str) -> None:
    if not address:
        return
    for shape in iter_shapes(slide.shapes):
        if not shape.has_table:
            continue
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text
                if "Ocean Blvd" in cell_text or "Fort Lauderdale" in cell_text or "2821" in cell_text:
                    set_cell_text(cell, address)
                    return


def replace_slide1_title3(slide, proposal_date: str, lob: str) -> None:
    a_rPr = qn("a:rPr")
    a_t = qn("a:t")
    targets = []
    if proposal_date:
        targets.append(("March 21, 2026", proposal_date.upper()))
    if lob:
        targets.append(("Commercial General Liability", lob.upper()))
    if not targets:
        return

    for shape in iter_shapes(slide.shapes):
        if shape.name != "Title 3" or not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text or ""
                if not text:
                    continue
                applicable = [(p, v) for p, v in targets if p in text]
                if not applicable:
                    continue

                original_rPr = run._r.find(a_rPr)
                preserved_rPr = deepcopy(original_rPr) if original_rPr is not None else None

                new_text = text
                for placeholder, value in applicable:
                    new_text = new_text.replace(placeholder, value)

                t_elem = run._r.find(a_t)
                if t_elem is not None:
                    t_elem.text = new_text
                else:
                    run.text = new_text

                if preserved_rPr is not None:
                    current_rPr = run._r.find(a_rPr)
                    if current_rPr is not None:
                        run._r.remove(current_rPr)
                    run._r.insert(0, preserved_rPr)
        return


def replace_client_name_on_slide1(slide, client_name: str) -> None:
    if not client_name:
        return
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join(run.text or "" for run in paragraph.runs)
            if "SUSAN LOMAN" not in full_text:
                continue
            new_text = full_text.replace("SUSAN LOMAN", client_name)
            if not paragraph.runs:
                continue
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""


# ---------------------------------------------------------------------------
# Output file naming (unchanged from original script)
# ---------------------------------------------------------------------------
def find_unique_output_path(base_name: str, output_dir: Path) -> Path:
    candidate = output_dir / base_name
    if not candidate.exists():
        return candidate
    stem = candidate.stem
    suffix = candidate.suffix
    index = 2
    while True:
        alternate = output_dir / f"{stem}_{index}{suffix}"
        if not alternate.exists():
            return alternate
        index += 1


def build_output_filename(proposal_date: Any, client_name: str) -> str:
    date_part = format_date_for_filename(proposal_date)
    client_part = sanitize_filename_part(client_name)
    return f"{date_part}_{client_part}.pptx"


def existing_quotes_for_row(row: pd.Series, output_dir: Path) -> List[Path]:
    """Return already-generated quote files for this row (used by the GUI
    to warn about duplicate generation)."""
    client_name = normalize_value(row["Client Name"])
    base = build_output_filename(row["Proposal Date"], client_name)
    stem = Path(base).stem
    if not output_dir.exists():
        return []
    return sorted(
        p for p in output_dir.glob(f"{stem}*.pptx")
        if p.stem == stem or re.fullmatch(re.escape(stem) + r"_\d+", p.stem)
    )


# ---------------------------------------------------------------------------
# Single-row quote generation (the original main() loop body, extracted)
# ---------------------------------------------------------------------------
def _blip_candidates(shapes):
    """Return (shape, blip) pairs for every image-bearing shape: real
    pictures AND shapes that use a picture FILL (templates often do the
    latter for cover images)."""
    candidates = []
    for shape in iter_shapes(shapes):
        try:
            blip = None
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blip = shape._element.blipFill.find(qn("a:blip"))
            else:
                spPr = shape._element.find(qn("p:spPr"))
                if spPr is not None:
                    blip_fill = spPr.find(qn("a:blipFill"))
                    if blip_fill is not None:
                        blip = blip_fill.find(qn("a:blip"))
            if blip is not None and (shape.width or 0) and (shape.height or 0):
                candidates.append((shape, blip))
        except Exception:
            continue
    return candidates


def _swap_image(blip, image_path, part) -> None:
    """Point an existing image reference at a new image file. The frame's
    position and size are untouched, so the layout never changes."""
    _, rId = part.get_or_add_image_part(str(image_path))
    blip.set(qn("r:embed"), rId)


def _fetch_row_image(row: pd.Series, field: str):
    """Get an image referenced by a submission field.

    JotForm file-upload fields export as a LINK to the uploaded file, so
    the value is usually a URL — it's downloaded to a temp file. A local
    file path works too. Returns a Path, or None (missing field, bad URL,
    no internet, ...) — callers treat None as 'no image provided'.
    """
    raw = normalize_value(row.get(field))
    if not raw:
        return None
    # An upload field can hold several links; use the first one
    first = re.split(r"[\s;,]+", raw)[0].strip()
    if not first:
        return None
    if first.lower().startswith(("http://", "https://")):
        try:
            import tempfile
            import urllib.request
            from urllib.parse import urlparse
            ext = Path(urlparse(first).path).suffix.lower()
            if ext not in (".png", ".jpg", ".jpeg"):
                ext = ".png"
            request = urllib.request.Request(
                first, headers={"User-Agent": "Mozilla/5.0 (QuoteGenerator)"}
            )
            with urllib.request.urlopen(request, timeout=20) as resp:
                data = resp.read()
            tmp = Path(tempfile.gettempdir()) / f"quotegen_{abs(hash(first))}{ext}"
            tmp.write_bytes(data)
            return tmp
        except Exception as exc:
            _log_error(exc)  # offline / login-protected upload / dead link
            return None
    local = Path(first)
    return local if local.exists() else None


def replace_cover_photo(prs: Presentation, photo_path) -> None:
    """Swap the large cover image on page 1 for a user-chosen photo.

    The image may live on the slide itself, on its layout, or on the
    master (designers place cover art in any of the three) — all are
    searched, largest image wins.
    """
    if not photo_path:
        return
    slide = prs.slides[0]
    layout = slide.slide_layout
    owners = [
        (slide.shapes, slide.part),
        (layout.shapes, layout.part),
        (layout.slide_master.shapes, layout.slide_master.part),
    ]
    for shapes, part in owners:
        candidates = _blip_candidates(shapes)
        if candidates:
            _, blip = max(
                candidates,
                key=lambda t: (t[0].width or 0) * (t[0].height or 0),
            )
            _swap_image(blip, photo_path, part)
            return


# ---------------------------------------------------------------------------
# Carrier logos.
#
# The proposal templates contain a carrier LOGO IMAGE on top of the carrier
# cells. To make it match the quoted carrier, drop logo images into a folder
# named "CarrierLogos" (next to the app/exe, or inside the Desktop "Quotes"
# folder). Name each file after the carrier: Chaucer.png, Lloyds.png,
# Citizens.jpg ... matching ignores case, spaces and punctuation.
# If no logo file matches the quoted carrier, the template's old logo is
# REMOVED so the carrier name (filled into the cell) shows instead of a
# wrong logo.
# ---------------------------------------------------------------------------
LOGO_DIR_NAME = "CarrierLogos"


def _carrier_key(name: str) -> str:
    return re.sub(r"[^a-z0-9]", "", str(name).lower())


def _find_carrier_logo(carrier_name: str):
    key = _carrier_key(carrier_name)
    if not key:
        return None
    search_dirs = [app_base_dir() / LOGO_DIR_NAME,
                   default_output_dir() / LOGO_DIR_NAME]
    bundle_dir = getattr(sys, "_MEIPASS", None)
    if bundle_dir:
        search_dirs.append(Path(bundle_dir) / LOGO_DIR_NAME)
    for folder in search_dirs:
        try:
            if not folder.exists():
                continue
            for p in sorted(folder.iterdir()):
                if p.suffix.lower() in (".png", ".jpg", ".jpeg") \
                        and _carrier_key(p.stem) == key:
                    return p
        except Exception:
            continue
    return None


def _update_carrier_logos(row: pd.Series, coverage_frames) -> None:
    """Swap (or remove) the logo image sitting on each coverage table's
    carrier cell. coverage_frames is [(slide, table_shape, prefix), ...]."""
    for slide, frame, prefix in coverage_frames:
        field = "Expiring Carrier" if prefix == "Current" else "New Carrier"
        carrier = normalize_value(row.get(field))
        try:
            fx, fy = frame.left, frame.top
            fw, fh = frame.width, frame.height
        except Exception:
            continue
        if None in (fx, fy, fw, fh):
            continue
        for shape in list(slide.shapes):
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                if None in (shape.left, shape.top, shape.width, shape.height):
                    continue
                cx = shape.left + shape.width / 2
                cy = shape.top + shape.height / 2
            except Exception:
                continue
            # A picture whose centre sits inside this table = carrier logo
            if not (fx <= cx <= fx + fw and fy <= cy <= fy + fh):
                continue
            # Priority: logo uploaded on the JotForm ("New Carrier Logo" /
            # "Expiring Carrier Logo") -> CarrierLogos folder -> remove
            logo = _fetch_row_image(row, field + " Logo")
            if logo is None:
                logo = _find_carrier_logo(carrier)
            blip = shape._element.blipFill.find(qn("a:blip"))
            if logo is not None and blip is not None:
                _swap_image(blip, logo, slide.part)
            elif carrier:
                # No logo available: remove the wrong logo so the carrier
                # name written into the cell is visible instead
                shape._element.getparent().remove(shape._element)


def generate_quote(
    row: pd.Series,
    template_path: Path | None = None,
    output_dir: Path | None = None,
    profile_key: str = "AUTO",
    photo_path: Path | None = None,
) -> Path:
    """Generate one proposal .pptx from a single submission row.

    AUTO uses the original hand-tuned slide-filling path (unchanged).
    Other quote types (HO3, ...) use the generic mapping-file-driven path.
    Returns the path of the saved file; raises QuoteGenerationError on
    failure.
    """
    template_path = Path(template_path) if template_path else default_template_path(profile_key)
    output_dir = Path(output_dir) if output_dir else default_output_dir()
    output_dir.mkdir(parents=True, exist_ok=True)  # ensure output folder exists

    if not template_path.exists():
        raise QuoteGenerationError(
            "The proposal template could not be found.\n\n"
            f"Please place the template file ('{template_path.name}') in the "
            "same folder as this application, or contact your administrator."
        )

    # No photo chosen in the app? Use the one uploaded on the JotForm, if any
    if photo_path is None:
        photo_path = _fetch_row_image(row, "Cover Photo")

    if profile_key != "AUTO":
        return _generate_generic(row, template_path, output_dir, profile_key, photo_path)

    try:
        mapping, named_insured_lines, client_name = build_text_mapping(row)
        prs = Presentation(template_path)
        replace_cover_photo(prs, photo_path)

        # Walk the first SLIDE_REPLACE_LIMIT slides and apply replacements
        for slide_index, slide in enumerate(prs.slides, start=1):
            if slide_index > SLIDE_REPLACE_LIMIT:
                break

            if slide_index == 1:
                # Slide 1 needs formatting-preserving replacement of the
                # client name, date and LOB, so those keys are excluded from
                # the generic mapping pass.
                replace_client_name_on_slide1(slide, client_name)
                replace_slide1_title3(
                    slide,
                    normalize_value(row["Proposal Date"]),
                    normalize_value(row["LOB"]),
                )
                excluded = {
                    "SUSAN LOMAN",
                    "TEST CO #2",
                    "March 21, 2026",
                    "Commercial General Liability",
                }
                slide_mapping = {k: v for k, v in mapping.items() if k not in excluded}
                replace_text_on_slide(slide, slide_mapping)
                continue

            if slide_index == 2:
                fill_slide2_named_insureds_textbox(slide, named_insured_lines)

            if slide_index == 3:
                fill_driver_table(slide, row)
                fill_vehicle_table(slide, row)
                fill_slide3_address_table(
                    slide, mapping["2821 N Ocean Blvd Fort Lauderdale Florida 33308"]
                )

            replace_text_on_slide(slide, mapping)

        output_filename = build_output_filename(row["Proposal Date"], client_name)
        output_path = find_unique_output_path(output_filename, output_dir)
        prs.save(output_path)
        return output_path
    except QuoteGenerationError:
        raise
    except Exception as exc:
        # Never surface raw Python errors to staff. Details go to a log file
        # the administrator can check; the user sees a friendly message.
        _log_error(exc)
        raise QuoteGenerationError(
            "The quote could not be created. The submission data may be "
            "incomplete or in an unexpected format.\n\n"
            "Please try another submission, or contact your administrator."
        ) from exc


# ---------------------------------------------------------------------------
# Generic generation path for additional quote types (HO3, ...).
#
# Driven by a mapping JSON (e.g. ho3_mapping.json) that pairs exact text
# strings found in the template with either a canonical field name or a
# special token. This keeps new quote types data-driven: update the JSON,
# no code changes needed.
#
# Special tokens usable as mapping values:
#   __client__            client name as typed     __client_upper__   UPPERCASE
#   __date__              proposal date            __date_upper__     UPPERCASE
#   __producer_upper__    "LEAD / SECONDARY" producers, uppercase
#   __address__           street, city, state, zip joined
#   __named_insureds__    names, one per line
#   __lob_upper__         line of business, uppercase
# ---------------------------------------------------------------------------
def _mapping_json_path(profile_key: str) -> Path | None:
    import quote_profiles
    name = quote_profiles.PROFILES[profile_key].get("mapping_json")
    if not name:
        return None
    beside = app_base_dir() / name
    if beside.exists():
        return beside
    bundle_dir = getattr(sys, "_MEIPASS", None)
    if bundle_dir and (Path(bundle_dir) / name).exists():
        return Path(bundle_dir) / name
    return beside  # not found; caller reports it


def _token_value(token: str, row: pd.Series) -> str:
    client = normalize_value(row.get("Client Name"))
    date = normalize_value(row.get("Proposal Date"))
    if token == "__client__":
        return client
    if token == "__client_upper__":
        return client.upper()
    if token == "__date__":
        return date
    if token == "__date_upper__":
        return date.upper()
    if token == "__producer_upper__":
        lead = normalize_value(row.get("Lead Producer"))
        second = normalize_value(row.get("Secondary Producer"))
        text = f"{lead} / {second}" if lead and second else (lead or second)
        return text.upper()
    if token == "__address__":
        return ", ".join(
            p for p in (
                normalize_value(row.get("Street")),
                normalize_value(row.get("City")),
                normalize_value(row.get("State")),
                normalize_value(row.get("ZIP")),
            ) if p
        )
    if token == "__named_insureds__":
        return "\n".join(get_named_insured_lines(normalize_value(row.get("Named Insureds"))))
    if token == "__lob_upper__":
        return normalize_value(row.get("LOB")).upper()
    if token in ("__ho3_total_fire_wind__", "__ho3_total_all__"):
        # Option totals on the HO3 options page: base premium (+ wind)(+ flood)
        def _money(field: str):
            digits = re.sub(r"[^0-9.]", "", normalize_value(row.get(field)))
            try:
                return float(digits) if digits else None
            except ValueError:
                return None
        parts = [_money("New Premium"), _money("Wind Premium")]
        if token == "__ho3_total_all__":
            parts.append(_money("Flood Premium"))
        if any(p is None for p in parts):
            return ""  # an option wasn't quoted; blank beats a wrong total
        return "$ {:,.0f}".format(sum(parts))
    return ""


# ---------------------------------------------------------------------------
# HO3 structural table filling.
#
# The HO3 template's tables reuse identical sample values for different
# fields (e.g. "$ 153,700" is both Other Structures and Loss of Use, and
# "Chaucer" is both the prior and proposed carrier), so those cells are
# located by their column headers and filled positionally — exactly like
# the Auto template's driver/vehicle tables. Plain text replacement then
# handles everything outside these tables.
# ---------------------------------------------------------------------------

# header keyword (lowercase) -> canonical field suffix
_HO3_COVERAGE_COLUMNS = [
    ("dwelling", "Dwelling Limit"),
    ("other", "Other Structures Limit"),
    ("personal property", "Personal Property Limit"),
    ("loss of use", "Loss of Use Limit"),
    ("liability", "Personal Liability Limit"),
    ("medical", "Medical Payments Limit"),
    ("aop", "AOP Deductible"),
    ("wind", "Hurricane Deductible"),
]


def _find_header_row(table, *keywords: str):
    """Index of the first row whose combined text contains all keywords."""
    for i, table_row in enumerate(table.rows):
        text = " ".join(c.text for c in table_row.cells).lower()
        if all(k in text for k in keywords):
            return i
    return None


def _fill_ho3_coverage_table(table, row: pd.Series, prefix: str) -> None:
    """Fill one coverage table ('Current' or 'New' fields by column header)."""
    header_idx = _find_header_row(table, "dwelling", "loss of use")
    if header_idx is None or header_idx + 1 >= len(table.rows):
        return
    headers = [c.text.lower() for c in table.rows[header_idx].cells]
    value_cells = table.rows[header_idx + 1].cells
    n = len(value_cells)
    for j in range(n):
        header = headers[j] if j < len(headers) else ""
        field = None
        for keyword, suffix in _HO3_COVERAGE_COLUMNS:
            if keyword in header:
                field = f"{prefix} {suffix}"
                break
        if field is None:
            # Unlabelled columns: first = carrier, last = annual premium
            if j == 0:
                field = "Expiring Carrier" if prefix == "Current" else "New Carrier"
            elif j == n - 1:
                field = "Current Annual Premium" if prefix == "Current" else "New Premium"
            else:
                continue
        # Always overwrite — a blank beats a leftover sample value
        set_cell_text(value_cells[j], normalize_value(row.get(field)))


def _fill_ho3_premium_summary(table, row: pd.Series) -> None:
    """Premium summary row: effective date, prior/proposed carrier+premium."""
    header_idx = _find_header_row(table, "expiring premium")
    if header_idx is None or header_idx + 1 >= len(table.rows):
        return
    headers = [c.text.lower() for c in table.rows[header_idx].cells]
    value_cells = table.rows[header_idx + 1].cells
    column_fields = [
        ("effective", "Policy Date"),
        ("prior", "Expiring Carrier"),
        ("expiring premium", "Current Annual Premium"),
        ("proposed carrier", "New Carrier"),
        ("proposed premium", "New Premium"),
    ]
    for j, header in enumerate(headers):
        if j >= len(value_cells):
            break
        for keyword, field in column_fields:
            if keyword in header:
                set_cell_text(value_cells[j], normalize_value(row.get(field)))
                break


def _fill_ho3_location_table(table, row: pd.Series) -> None:
    """Property address row (kept structural: replacing the text 'FL' would
    also corrupt the 'FL License #...' footer on every slide)."""
    header_idx = _find_header_row(table, "address", "zip")
    if header_idx is None or header_idx + 1 >= len(table.rows):
        return
    headers = [c.text.lower() for c in table.rows[header_idx].cells]
    value_cells = table.rows[header_idx + 1].cells
    column_fields = [
        ("address", "Street"),
        ("city", "City"),
        ("state", "State"),
        ("zip", "ZIP"),
    ]
    for j, header in enumerate(headers):
        if j >= len(value_cells):
            break
        for keyword, field in column_fields:
            if keyword in header:
                set_cell_text(value_cells[j], normalize_value(row.get(field)))
                break


def _fill_ho3_tables(prs: Presentation, row: pd.Series) -> None:
    """Classify and fill every structured table in the HO3 template.

    Coverage tables are matched in document order: the first is the
    current/expiring policy, the second is the proposed one (matching the
    template layout: 'Current Coverage Summary' above 'Option A').
    """
    coverage_seen = 0
    coverage_frames = []  # (slide, table shape, "Current"/"New")
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if not shape.has_table:
                continue
            table = shape.table
            if _find_header_row(table, "expiring premium") is not None:
                _fill_ho3_premium_summary(table, row)
            elif _find_header_row(table, "dwelling", "loss of use") is not None:
                prefix = "Current" if coverage_seen == 0 else "New"
                _fill_ho3_coverage_table(table, row, prefix)
                coverage_frames.append((slide, shape, prefix))
                coverage_seen += 1
            elif _find_header_row(table, "address", "zip") is not None \
                    and _find_header_row(table, "entity") is None:
                # Locations table (the Additional Insured table also has
                # address columns but starts with 'Entity' — leave it alone)
                _fill_ho3_location_table(table, row)

    # Match (or remove) the carrier logo images sitting on the tables
    _update_carrier_logos(row, coverage_frames)


_STRUCTURAL_FILLERS = {
    "ho3": _fill_ho3_tables,
}


def _generate_generic(row: pd.Series, template_path: Path,
                      output_dir: Path, profile_key: str,
                      photo_path: Path | None = None) -> Path:
    """Mapping-file-driven generation for non-AUTO quote types."""
    mapping_path = _mapping_json_path(profile_key)
    if mapping_path is None or not mapping_path.exists():
        raise QuoteGenerationError(
            "This quote type isn't fully set up yet.\n\n"
            "Please contact your administrator. "
            f"(Missing mapping file: {mapping_path.name if mapping_path else profile_key})"
        )

    try:
        import json
        with open(mapping_path, encoding="utf-8") as fh:
            raw_mapping = json.load(fh)

        # Build placeholder -> replacement text. Keys starting with "_"
        # (like "_README") are documentation, not placeholders.
        mapping = {}
        for placeholder, source in raw_mapping.items():
            if placeholder.startswith("_"):
                continue
            if isinstance(source, str) and source.startswith("__"):
                mapping[placeholder] = _token_value(source, row)
            else:
                mapping[placeholder] = normalize_value(row.get(source))

        client_name = normalize_value(row.get("Client Name"))
        prs = Presentation(template_path)
        replace_cover_photo(prs, photo_path)

        # 1) Structural, table-aware filling first (clears ambiguous sample
        #    values so the text pass below can't mis-replace them)
        import quote_profiles
        filler_key = quote_profiles.PROFILES[profile_key].get("structural")
        filler = _STRUCTURAL_FILLERS.get(filler_key)
        if filler:
            filler(prs, row)

        # 2) Then plain text replacement everywhere else
        for slide in prs.slides:
            replace_text_on_slide(slide, mapping)

        output_filename = build_output_filename(row.get("Proposal Date"), client_name)
        output_path = find_unique_output_path(output_filename, output_dir)
        prs.save(output_path)
        return output_path
    except QuoteGenerationError:
        raise
    except Exception as exc:
        _log_error(exc)
        raise QuoteGenerationError(
            "The quote could not be created. The submission data may be "
            "incomplete or in an unexpected format.\n\n"
            "Please try another submission, or contact your administrator."
        ) from exc


def _log_error(exc: Exception) -> None:
    """Append technical error details to Quotes/error_log.txt for debugging."""
    try:
        import datetime
        import traceback
        log_path = default_output_dir() / "error_log.txt"
        with open(log_path, "a", encoding="utf-8") as fh:
            fh.write(f"\n--- {datetime.datetime.now():%Y-%m-%d %H:%M:%S} ---\n")
            fh.write("".join(traceback.format_exception(type(exc), exc, exc.__traceback__)))
    except Exception:
        pass  # logging must never crash the app
